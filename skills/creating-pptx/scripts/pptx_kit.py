"""
pptx_kit.py
DualSlide: python-pptx と Pillow を1回の呼び出しで同時描画するヘルパー。
コンサル品質の 16:9 (13.333×7.5 in) スライドを .pptx と確認用 .png に同時生成する。

依存ライブラリ:
    python-pptx >= 1.0
    Pillow       >= 10.0
    lxml         >= 4.9  (python-pptx の依存として自動インストール)

使用例:
    ds = DualSlide()
    ds.head_message("コスト削減余地は年間 3 億円")
    ds.rect(0.55, 1.2, 12.2, 5.5)
    ds.save("output.pptx", "output.png")
"""
from __future__ import annotations

import math
import os
import subprocess
import warnings
from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

from PIL import Image, ImageDraw, ImageFont

# ---------------------------------------------------------------------------
# カラーパレット（この 6 色以外はモジュール全体で使用禁止）
# ---------------------------------------------------------------------------
BLUE   = (0x19, 0x76, 0xD2)   # #1976D2  タイトル/見出し帯/バッジ/強調系列
VERMIL = (0xD8, 0x43, 0x15)   # #D84315  強調キーワード専用（ベタ塗り面に使わない）
LIGHT  = (0xE8, 0xF4, 0xFB)   # #E8F4FB  淡い塗り（最も強調したい 1 ブロックのみ）
INK    = (0x22, 0x22, 0x22)   # #222222  本文（純黒禁止）
SUB    = (0x88, 0x88, 0x88)   # #888888  補助・脚注・出典
LINE   = (0xE0, 0xE0, 0xE0)   # #E0E0E0  枠線・ヘアライン・破線
WHITE  = (0xFF, 0xFF, 0xFF)   # #FFFFFF  背景

# スライド定数
SLIDE_W_IN: float = 13.333
SLIDE_H_IN: float = 7.5
PPTX_FONT: str = "Noto Sans JP"

# 型エイリアス
ColorTuple = tuple[int, int, int]
RunItem = tuple[str, dict]
RunsType = Union[str, list[RunItem]]

# stack100 / waterfall で使用するセグメント色シーケンス（パレット内で循環）
_SEG_COLORS: list[ColorTuple] = [BLUE, SUB, LINE, LIGHT]

# PIL フォントキャッシュ（size_pt × bold のペアでキャッシュ）
_FONT_CACHE: dict[tuple[float, bool], "ImageFont.FreeTypeFont"] = {}


# ---------------------------------------------------------------------------
# モジュールレベルのユーティリティ関数
# ---------------------------------------------------------------------------

def _resolve_font(size_pt: float, bold: bool = False) -> "ImageFont.FreeTypeFont":
    """
    PIL 用の日本語フォントを解決する。
    優先順: fc-match → 既知パス → デフォルトフォント（警告付き）。
    """
    key = (size_pt, bold)
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]

    size_px = max(8, int(size_pt * 150 / 72))
    style = "Bold" if bold else "Regular"

    # fc-match で最適パスを取得
    try:
        result = subprocess.run(
            ["fc-match", "--format=%{file}", f"Noto Sans CJK JP:style={style}"],
            capture_output=True, text=True, timeout=2,
        )
        if result.returncode == 0:
            p = result.stdout.strip()
            if p and os.path.exists(p):
                font = ImageFont.truetype(p, size_px)
                _FONT_CACHE[key] = font
                return font
    except Exception:
        pass

    # 既知パスを順に試す
    candidates: list[tuple[str, int | None]] = [
        # macOS ヒラギノ角ゴシック
        (
            "/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc" if bold
            else "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
            0,
        ),
        # Linux Noto Sans CJK
        ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 0),
        ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", 0),
        ("/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc", 0),
        ("/usr/local/share/fonts/NotoSansJP-Regular.ttf", None),
        ("/usr/share/fonts/truetype/noto/NotoSansJP-Regular.ttf", None),
    ]
    for path, idx in candidates:
        if os.path.exists(path):
            try:
                font = (
                    ImageFont.truetype(path, size_px, index=idx)
                    if idx is not None
                    else ImageFont.truetype(path, size_px)
                )
                _FONT_CACHE[key] = font
                return font
            except Exception:
                pass

    warnings.warn(
        f"日本語フォントが未発見 (size={size_pt}pt, bold={bold})。"
        "豆腐(□)が表示される可能性があります。デフォルトフォントで継続します。",
        stacklevel=4,
    )
    try:
        font = ImageFont.load_default(size=size_px)
    except TypeError:
        font = ImageFont.load_default()
    _FONT_CACHE[key] = font
    return font


def _suppress_effects(shape) -> None:
    """影・グラデ・3D 効果を XML レベルで無効化する（全図形に適用）。"""
    try:
        sp_el = shape._element
        spPr = sp_el.find(qn("p:spPr"))
        if spPr is None:
            return
        for tag in ("a:effectLst", "a:effectDag", "a:scene3d", "a:sp3d"):
            for el in list(spPr.findall(qn(tag))):
                spPr.remove(el)
        etree.SubElement(spPr, qn("a:effectLst"))
    except Exception:
        pass


def _set_run_font(
    run, name: str, size_pt: float, bold: bool, color: ColorTuple
) -> None:
    """pptx の run にフォント・サイズ・色・東アジアフォントを設定する。"""
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    # 日本語グリフを Noto Sans JP で描画するため a:ea を設定
    try:
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", name)
    except Exception:
        pass


def _set_no_line(shape_element) -> None:
    """spPr に <a:ln><a:noFill/></a:ln> を追加して枠線を消す。"""
    try:
        spPr = shape_element.find(qn("p:spPr"))
        if spPr is None:
            return
        ln_el = spPr.find(qn("a:ln"))
        if ln_el is None:
            ln_el = etree.SubElement(spPr, qn("a:ln"))
        # 既存のsolidFill等を削除してnoFillに
        for child in list(ln_el):
            ln_el.remove(child)
        etree.SubElement(ln_el, qn("a:noFill"))
    except Exception:
        pass


def _normalize_runs(
    runs: RunsType,
    default_size: float,
    default_bold: bool,
    default_color: ColorTuple,
) -> list[RunItem]:
    """runs を list[RunItem] に正規化する。"""
    if isinstance(runs, str):
        return [(runs, {"size": default_size, "bold": default_bold, "color": default_color})]
    result: list[RunItem] = []
    for item in runs:
        text, style = item
        result.append((
            text,
            {
                "size":  style.get("size",  default_size),
                "bold":  style.get("bold",  default_bold),
                "color": style.get("color", default_color),
            },
        ))
    return result


def _break_runs_to_lines(
    runs: list[RunItem],
    draw: ImageDraw.ImageDraw,
    max_w_px: float,
) -> list[list[RunItem]]:
    """
    ミックスrunを行単位に分解する。
    文字単位で折り返す（日本語・英語混在対応）。
    """
    lines: list[list[RunItem]] = []
    cur_line: list[RunItem] = []
    cur_w = 0.0

    for text, style in runs:
        font = _resolve_font(style["size"], style["bold"])
        buf = ""
        for char in text:
            if char == "\n":
                if buf:
                    cur_line.append((buf, style))
                    buf = ""
                lines.append(cur_line)
                cur_line = []
                cur_w = 0.0
                continue
            cw = draw.textlength(char, font=font)
            if cur_w + cw > max_w_px and (cur_line or buf):
                if buf:
                    cur_line.append((buf, style))
                    buf = ""
                lines.append(cur_line)
                cur_line = []
                cur_w = 0.0
            buf += char
            cur_w += cw
        if buf:
            cur_line.append((buf, style))

    if cur_line:
        lines.append(cur_line)
    return lines


def _draw_dash_line(
    draw: ImageDraw.ImageDraw,
    x1: float, y1: float,
    x2: float, y2: float,
    color: ColorTuple,
    width: int,
    dash: int = 6,
    gap: int = 4,
) -> None:
    """PIL でドット/ダッシュ線を描画する。"""
    dx = x2 - x1
    dy = y2 - y1
    length = math.hypot(dx, dy)
    if length < 1:
        return
    ux = dx / length
    uy = dy / length
    pos = 0.0
    drawing = True
    while pos < length:
        end_pos = min(pos + (dash if drawing else gap), length)
        if drawing:
            draw.line(
                [
                    (x1 + ux * pos, y1 + uy * pos),
                    (x1 + ux * end_pos, y1 + uy * end_pos),
                ],
                fill=color,
                width=width,
            )
        pos = end_pos
        drawing = not drawing


def _add_dash_to_connector(cnx) -> None:
    """pptx コネクタの line に sysDot ダッシュスタイルを設定する。"""
    try:
        ln = cnx.line._ln
        if ln is not None:
            prstDash = ln.find(qn("a:prstDash"))
            if prstDash is None:
                prstDash = etree.SubElement(ln, qn("a:prstDash"))
            prstDash.set("val", "sysDot")
    except Exception:
        pass


# ---------------------------------------------------------------------------
# DualSlide クラス
# ---------------------------------------------------------------------------

class DualSlide:
    """
    1回の呼び出しで python-pptx スライドと Pillow プレビュー画像の
    両方へ同時に描画するキャンバス。

    座標系: インチ、原点は左上。
    プレビューサイズ: preview_scale px/inch (デフォルト 150 → 2000×1125 px)。

    Usage::

        ds = DualSlide()
        ds.head_message("So What が入る結論一文")
        ds.rect(0.55, 1.2, 12.2, 5.5)
        ds.source("出所: 社内調査 (2024年)")
        ds.save("result.pptx", "result.png")
    """

    def __init__(self, preview_scale: int = 150) -> None:
        self._scale = preview_scale

        # python-pptx 側の初期化
        self._prs = Presentation()
        self._prs.slide_width  = Inches(SLIDE_W_IN)
        self._prs.slide_height = Inches(SLIDE_H_IN)
        blank_layout = self._prs.slide_layouts[6]  # 空白レイアウト
        self._slide = self._prs.slides.add_slide(blank_layout)

        # Pillow 側の初期化
        w_px = int(SLIDE_W_IN * preview_scale)
        h_px = int(SLIDE_H_IN * preview_scale)
        self._img  = Image.new("RGB", (w_px, h_px), WHITE)
        self._draw = ImageDraw.Draw(self._img)

    # ------------------------------------------------------------------
    # 内部単位変換ヘルパー
    # ------------------------------------------------------------------

    def _px(self, inch: float) -> float:
        """インチ → ピクセル（float）"""
        return inch * self._scale

    def _ptpx(self, pt: float) -> int:
        """ポイント → ピクセル（int）"""
        return max(1, int(pt * self._scale / 72))

    # ------------------------------------------------------------------
    # 公開 API
    # ------------------------------------------------------------------

    def head_message(self, text: str, emphasis: str | None = None) -> None:
        """
        スライド上部に結論一文を配置する（HEAD 19pt、INK）。
        emphasis に指定した部分文字列のみ VERMIL かつ Bold で描画する。
        """
        m = 0.55
        x, y, w, h = m, m, SLIDE_W_IN - 2 * m, 0.55

        if emphasis is None or emphasis not in text:
            runs: list[RunItem] = [(text, {"size": 19.0, "bold": False, "color": INK})]
        else:
            idx = text.index(emphasis)
            runs = []
            if idx > 0:
                runs.append((text[:idx], {"size": 19.0, "bold": False, "color": INK}))
            runs.append((emphasis, {"size": 19.0, "bold": True, "color": VERMIL}))
            after = text[idx + len(emphasis):]
            if after:
                runs.append((after, {"size": 19.0, "bold": False, "color": INK}))

        self.text(x, y, w, h, runs, size=19.0)

    def text(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        runs: RunsType,
        size: float = 12.0,
        color: ColorTuple = INK,
        bold: bool = False,
        align: str = "left",
        valign: str = "top",
        line_spacing: float = 1.4,
    ) -> None:
        """
        テキストボックスを描画する。

        runs: str または [(text, {color, bold, size}), ...] で指定。
        インライン強調（VERMIL 等）は後者の形式で指定する。
        word_wrap が有効。座標はすべてインチ。
        """
        norm = _normalize_runs(runs, size, bold, color)

        # --- pptx 側 ---
        txBox = self._slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        try:
            from pptx.enum.text import MSO_ANCHOR
            vmap = {
                "top":    MSO_ANCHOR.TOP,
                "middle": MSO_ANCHOR.MIDDLE,
                "bottom": MSO_ANCHOR.BOTTOM,
            }
            tf.vertical_anchor = vmap.get(valign, MSO_ANCHOR.TOP)
        except Exception:
            pass

        amap = {
            "left":   PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right":  PP_ALIGN.RIGHT,
        }

        # 改行文字で段落に分割
        paras: list[list[RunItem]] = [[]]
        for rtext, rstyle in norm:
            parts = rtext.split("\n")
            for i, part in enumerate(parts):
                if i > 0:
                    paras.append([])
                if part:
                    paras[-1].append((part, rstyle))

        for pi, para_runs in enumerate(paras):
            para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
            para.alignment = amap.get(align, PP_ALIGN.LEFT)
            para.line_spacing = line_spacing
            for rtext, rstyle in para_runs:
                run = para.add_run()
                run.text = rtext
                _set_run_font(run, PPTX_FONT, rstyle["size"], rstyle["bold"], rstyle["color"])

        _suppress_effects(txBox)

        # --- PIL 側 ---
        x_px = self._px(x)
        y_px = self._px(y)
        w_px = self._px(w)
        h_px = self._px(h)

        pil_lines = _break_runs_to_lines(norm, self._draw, w_px)
        if not pil_lines:
            return

        def _line_height(ln: list[RunItem]) -> float:
            if not ln:
                return self._ptpx(size) * 1.0
            heights = []
            for _, rs in ln:
                fnt = _resolve_font(rs["size"], rs["bold"])
                bb = self._draw.textbbox((0, 0), "Ag", font=fnt)
                heights.append(float(bb[3] - bb[1]))
            return max(heights)

        line_hs = [_line_height(ln) for ln in pil_lines]
        total_h = sum(lh * line_spacing for lh in line_hs)

        if valign == "middle":
            y_cur = y_px + (h_px - total_h) / 2
        elif valign == "bottom":
            y_cur = y_px + h_px - total_h
        else:
            y_cur = y_px

        for ln, lh in zip(pil_lines, line_hs):
            lw = sum(
                self._draw.textlength(rt, font=_resolve_font(rs["size"], rs["bold"]))
                for rt, rs in ln
            )
            if align == "center":
                x_cur = x_px + (w_px - lw) / 2
            elif align == "right":
                x_cur = x_px + w_px - lw
            else:
                x_cur = x_px

            for rt, rs in ln:
                fnt = _resolve_font(rs["size"], rs["bold"])
                self._draw.text((x_cur, y_cur), rt, font=fnt, fill=rs["color"])
                x_cur += self._draw.textlength(rt, font=fnt)

            y_cur += lh * line_spacing

    def rect(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: ColorTuple | None = None,
        line: ColorTuple | None = LINE,
        line_w: float = 1.0,
        radius: float = 0.0,
    ) -> None:
        """
        矩形を描画する（フラット・影なし）。
        fill=None で塗りなし、line=None で枠線なし。
        radius > 0 で角丸（全角統一）。
        """
        # --- pptx 側 ---
        shape_type = 5 if radius > 0 else 1   # 5=ROUNDED_RECTANGLE, 1=RECTANGLE
        shape = self._slide.shapes.add_shape(
            shape_type,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        # 角丸半径の設定
        if radius > 0:
            try:
                half_min = min(w, h) / 2
                adj_val = int(min(50000, (radius / half_min) * 50000))
                spPr = shape._element.find(qn("p:spPr"))
                if spPr is not None:
                    prstGeom = spPr.find(qn("a:prstGeom"))
                    if prstGeom is not None:
                        avLst = prstGeom.find(qn("a:avLst"))
                        if avLst is None:
                            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
                        gd = avLst.find(qn("a:gd"))
                        if gd is None:
                            gd = etree.SubElement(avLst, qn("a:gd"))
                        gd.set("name", "adj")
                        gd.set("fmla", f"val {adj_val}")
            except Exception:
                pass
        # 塗り
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill)
        # 枠線
        if line is None:
            _set_no_line(shape._element)
        else:
            shape.line.color.rgb = RGBColor(*line)
            shape.line.width = Pt(line_w)
        _suppress_effects(shape)

        # --- PIL 側 ---
        x1 = self._px(x)
        y1 = self._px(y)
        x2 = self._px(x + w)
        y2 = self._px(y + h)
        r_px = int(self._px(radius))
        lw_px = self._ptpx(line_w) if line else 0

        if radius > 0:
            self._draw.rounded_rectangle(
                [x1, y1, x2, y2],
                radius=r_px,
                fill=fill,
                outline=line,
                width=lw_px if line else 0,
            )
        else:
            self._draw.rectangle(
                [x1, y1, x2, y2],
                fill=fill,
                outline=line,
                width=lw_px if line else 0,
            )

    def hline(
        self,
        x: float,
        y: float,
        w: float,
        color: ColorTuple = LINE,
        weight: float = 1.0,
        dash: bool = False,
    ) -> None:
        """水平線を描画する。dash=True でドット線。"""
        # --- pptx 側 ---
        try:
            cnx = self._slide.shapes.add_connector(
                1,   # MSO_CONNECTOR.STRAIGHT
                Inches(x), Inches(y), Inches(x + w), Inches(y),
            )
            cnx.line.color.rgb = RGBColor(*color)
            cnx.line.width = Pt(weight)
            if dash:
                _add_dash_to_connector(cnx)
        except Exception:
            pass

        # --- PIL 側 ---
        x1_px = self._px(x)
        y_px  = self._px(y)
        x2_px = self._px(x + w)
        lw = self._ptpx(weight)
        if dash:
            _draw_dash_line(self._draw, x1_px, y_px, x2_px, y_px, color, lw)
        else:
            self._draw.line([(x1_px, y_px), (x2_px, y_px)], fill=color, width=lw)

    def vline(
        self,
        x: float,
        y: float,
        h: float,
        color: ColorTuple = LINE,
        weight: float = 1.0,
        dash: bool = False,
    ) -> None:
        """垂直線を描画する。dash=True でドット線。"""
        # --- pptx 側 ---
        try:
            cnx = self._slide.shapes.add_connector(
                1,   # MSO_CONNECTOR.STRAIGHT
                Inches(x), Inches(y), Inches(x), Inches(y + h),
            )
            cnx.line.color.rgb = RGBColor(*color)
            cnx.line.width = Pt(weight)
            if dash:
                _add_dash_to_connector(cnx)
        except Exception:
            pass

        # --- PIL 側 ---
        x_px  = self._px(x)
        y1_px = self._px(y)
        y2_px = self._px(y + h)
        lw = self._ptpx(weight)
        if dash:
            _draw_dash_line(self._draw, x_px, y1_px, x_px, y2_px, color, lw)
        else:
            self._draw.line([(x_px, y1_px), (x_px, y2_px)], fill=color, width=lw)

    def arrow(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        color: ColorTuple = BLUE,
        weight: float = 1.5,
    ) -> None:
        """直線矢印を描画する（x2,y2 側に矢じり・サイズ固定）。"""
        # --- pptx 側 ---
        try:
            cnx = self._slide.shapes.add_connector(
                1,   # MSO_CONNECTOR.STRAIGHT
                Inches(x1), Inches(y1), Inches(x2), Inches(y2),
            )
            cnx.line.color.rgb = RGBColor(*color)
            cnx.line.width = Pt(weight)
            try:
                ln = cnx.line._ln
                if ln is not None:
                    tailEnd = ln.find(qn("a:tailEnd"))
                    if tailEnd is None:
                        tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
                    tailEnd.set("type", "arrow")
                    tailEnd.set("w", "med")
                    tailEnd.set("len", "med")
            except Exception:
                pass
        except Exception:
            pass

        # --- PIL 側 ---
        x1_px = self._px(x1)
        y1_px = self._px(y1)
        x2_px = self._px(x2)
        y2_px = self._px(y2)
        lw = self._ptpx(weight)
        self._draw.line([(x1_px, y1_px), (x2_px, y2_px)], fill=color, width=lw)
        # 矢じり（等辺三角形）
        angle = math.atan2(y2_px - y1_px, x2_px - x1_px)
        alen = lw * 6
        aw   = lw * 3
        perp = angle + math.pi / 2
        bx = x2_px - alen * math.cos(angle)
        by = y2_px - alen * math.sin(angle)
        p1 = (bx + aw * math.cos(perp), by + aw * math.sin(perp))
        p2 = (bx - aw * math.cos(perp), by - aw * math.sin(perp))
        self._draw.polygon([(x2_px, y2_px), p1, p2], fill=color)

    def badge(
        self,
        x: float,
        y: float,
        n: int,
        d: float = 0.32,
    ) -> None:
        """
        丸数字バッジを描画する（BLUE 地・白数字・中央寄せ）。
        (x, y) はバッジ中心座標（インチ）。d は直径（インチ）。
        """
        # --- pptx 側 ---
        try:
            shape = self._slide.shapes.add_shape(
                9,   # OVAL
                Inches(x - d / 2), Inches(y - d / 2),
                Inches(d), Inches(d),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*BLUE)
            _set_no_line(shape._element)
            tf = shape.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(n)
            _set_run_font(run, PPTX_FONT, 9.0, True, WHITE)
            _suppress_effects(shape)
        except Exception:
            pass

        # --- PIL 側 ---
        cx = self._px(x)
        cy = self._px(y)
        r  = self._px(d / 2)
        self._draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=BLUE, outline=None)
        fnt = _resolve_font(9.0, True)
        txt = str(n)
        bb  = self._draw.textbbox((0, 0), txt, font=fnt)
        tw  = bb[2] - bb[0]
        th  = bb[3] - bb[1]
        self._draw.text((cx - tw / 2, cy - th / 2), txt, font=fnt, fill=WHITE)

    def table(
        self,
        x: float,
        y: float,
        w: float,
        rows: list[list[str]],
        col_widths: list[float],
        header: bool = True,
        aligns: list[str] | None = None,
        right_cols: set[int] | None = None,
    ) -> None:
        """
        表を描画する（横罫線のみ・縦罫線なし）。

        col_widths は相対幅（合計に占める比率）。
        header=True で先頭行を LIGHT 薄塗り + Bold。
        right_cols で指定した列インデックスを右揃えにする（数値列向け）。
        """
        if not rows or not col_widths:
            return
        total_cw = sum(col_widths)
        col_w_in = [w * cw / total_cw for cw in col_widths]
        row_h = 0.36   # 1行あたりの高さ（インチ）

        for ri, row in enumerate(rows):
            is_hdr = header and ri == 0
            ry = y + ri * row_h
            if is_hdr:
                self.rect(x, ry, w, row_h, fill=LIGHT, line=None)

            xc = x
            for ci, cell in enumerate(row):
                cw_in = col_w_in[ci] if ci < len(col_w_in) else 0
                if right_cols and ci in right_cols:
                    cell_align = "right"
                elif aligns and ci < len(aligns):
                    cell_align = aligns[ci]
                else:
                    cell_align = "left"
                self.text(
                    xc + 0.07, ry + 0.04,
                    cw_in - 0.14, row_h - 0.08,
                    cell,
                    size=10.5,
                    color=INK,
                    bold=is_hdr,
                    align=cell_align,
                    valign="middle",
                    line_spacing=1.2,
                )
                xc += cw_in
            self.hline(x, ry + row_h, w, color=LINE, weight=0.75)

        # ヘッダー上の罫線
        self.hline(x, y, w, color=LINE, weight=0.75)

    def bar(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        data: list[tuple[str, float]],
        highlight: int | None = None,
        unit: str = "",
        orient: str = "h",
    ) -> None:
        """
        横棒グラフを描画する（orient="h"）。
        highlight で指定したインデックスのバーを BLUE、他を SUB で描画。
        データラベルはバーの右に直接表示。
        """
        if not data:
            return
        max_val = max((v for _, v in data), default=1.0)
        if max_val <= 0:
            max_val = 1.0

        n = len(data)
        slot_h    = h / n
        bar_h     = slot_h * 0.62
        label_w   = min(w * 0.28, 2.0)
        val_w     = min(w * 0.12, 0.75)
        bar_area  = w - label_w - val_w - 0.06

        for i, (label, value) in enumerate(data):
            by = y + i * slot_h + (slot_h - bar_h) / 2
            bw = bar_area * (value / max_val)
            col = BLUE if i == highlight else SUB

            self.text(
                x, by, label_w - 0.06, bar_h,
                label, size=10.0, color=INK,
                align="right", valign="middle",
            )
            if bw > 0:
                self.rect(x + label_w, by, bw, bar_h, fill=col, line=None)
            self.text(
                x + label_w + bw + 0.06, by, val_w, bar_h,
                f"{value:g}{unit}",
                size=9.5, color=col, bold=(i == highlight),
                valign="middle",
            )

    def stack100(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        segments: list[tuple[str, float]],
    ) -> None:
        """
        100% 積み上げ横棒グラフを描画する。
        segments=[(label, value), ...] — value の合計を 100% とする。
        直接パーセントラベルを表示。
        """
        if not segments:
            return
        total = sum(v for _, v in segments) or 1.0
        bar_h = h * 0.55
        bar_y = y + (h - bar_h) / 2

        xc = x
        for i, (label, value) in enumerate(segments):
            sw = w * (value / total)
            col = _SEG_COLORS[i % len(_SEG_COLORS)]
            self.rect(xc, bar_y, sw, bar_h, fill=col, line=None)

            pct = value / total * 100
            if sw > 0.45:   # 幅が十分な場合のみ内側ラベル
                txt_col = WHITE if col in (BLUE, SUB) else INK
                self.text(
                    xc, bar_y, sw, bar_h,
                    f"{pct:.0f}%",
                    size=9.0, color=txt_col,
                    align="center", valign="middle",
                )
            self.text(
                xc, bar_y + bar_h + 0.05, sw, 0.25,
                label, size=8.5, color=SUB, align="center",
            )
            xc += sw

    def waterfall(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        steps: list[tuple[str, float, str]],
    ) -> None:
        """
        ウォーターフォールチャートを描画する。
        steps=[(label, value, kind)]
        kind: "start" | "inc" | "dec" | "total"
        増加 (inc) と start/total は BLUE、減少 (dec) は VERMIL で描画。
        """
        if not steps:
            return
        n = len(steps)
        slot_w = w / n
        bar_w  = slot_w * 0.62

        # ベース・高さを計算
        running = 0.0
        bars: list[tuple[str, float, str, float]] = []
        for label, value, kind in steps:
            if kind in ("start", "total"):
                bars.append((label, abs(value), kind, 0.0))
                running = value
            elif kind == "inc":
                bars.append((label, abs(value), kind, running))
                running += value
            else:   # dec
                running += value
                bars.append((label, abs(value), kind, running))

        all_tops = [base + val for _, val, _, base in bars]
        max_val = max(all_tops) if all_tops else 1.0
        if max_val <= 0:
            max_val = 1.0

        for i, (label, val, kind, base) in enumerate(bars):
            bx = x + i * slot_w + (slot_w - bar_w) / 2
            bh_scaled = h * (val / max_val)
            by = y + h - h * ((base + val) / max_val)
            col = VERMIL if kind == "dec" else BLUE

            if bh_scaled > 0:
                self.rect(bx, by, bar_w, bh_scaled, fill=col, line=None)

            sign = "-" if kind == "dec" else ("+" if kind == "inc" else "")
            self.text(
                bx, max(y, by - 0.28), bar_w, 0.26,
                f"{sign}{val:g}", size=9.0, color=col, align="center",
            )
            self.text(
                bx, y + h + 0.05, bar_w, 0.22,
                label, size=8.5, color=SUB, align="center",
            )

    def source(self, text: str) -> None:
        """スライド下部に出典テキストを配置する（SUB 8pt）。"""
        m = 0.55
        self.text(
            m, SLIDE_H_IN - m - 0.22,
            SLIDE_W_IN - 2 * m, 0.22,
            text,
            size=8.0, color=SUB,
        )

    def save(self, pptx_path: str, png_path: str) -> None:
        """
        pptx と png の両方を書き出す。
        出力先ディレクトリが存在しない場合は自動作成する。
        """
        Path(pptx_path).parent.mkdir(parents=True, exist_ok=True)
        Path(png_path).parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(pptx_path)
        self._img.save(png_path)
